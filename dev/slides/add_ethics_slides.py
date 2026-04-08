#!/usr/bin/env python3
"""
add_ethics_slides.py — Adds AI Ethics slides to the Cervantes Google Slides PPTX.

Inserts two slides (after slide 8 "Tech Highlights") covering:
  1. Responsible AI Framework
  2. Data Privacy & Student Protection

Uses the same La Mancha Parchment design system as generate_pptx.py.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Design tokens — La Mancha Parchment
# ---------------------------------------------------------------------------
TERRACOTTA = RGBColor(200, 90, 50)
WHEAT_GOLD = RGBColor(196, 162, 69)
PARCHMENT = RGBColor(244, 241, 234)
NEAR_BLACK = RGBColor(30, 28, 24)
WHITE = RGBColor(255, 255, 255)
BODY_TEXT = RGBColor(58, 55, 48)
MUTED_TEXT = RGBColor(107, 101, 90)
MASTERY = RGBColor(59, 130, 126)
MISCONCEPTION = RGBColor(212, 163, 71)
CRITICAL = RGBColor(158, 59, 59)
DARK_BG = RGBColor(42, 39, 32)

FONT_NAME = "Plus Jakarta Sans"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M_LEFT = Inches(0.8)
M_TOP = Inches(0.6)
CONTENT_W = Inches(11.7)

INPUT_FILE = Path(__file__).parent / "cervantes-google.pptx"
OUTPUT_FILE = Path(__file__).parent / "cervantes-deck-ethics.pptx"


# ---------------------------------------------------------------------------
# Helpers (same as generate_pptx.py)
# ---------------------------------------------------------------------------
def _font(run, size=18, bold=False, color=NEAR_BLACK, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _add_textbox(slide, left, top, width, height, text="", size=18,
                 bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _font(run, size=size, bold=bold, color=color)
    return txBox


def _add_para(text_frame, text, size=18, bold=False, color=NEAR_BLACK,
              space_before=Pt(0), space_after=Pt(4), align=PP_ALIGN.LEFT):
    p = text_frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    _font(run, size=size, bold=bold, color=color)
    return p


def _add_bullet(text_frame, text, size=18, color=BODY_TEXT,
                bold_prefix=""):
    p = text_frame.add_paragraph()
    p.space_before = Pt(2)
    p.space_after = Pt(4)
    if bold_prefix:
        r1 = p.add_run()
        r1.text = "\u2022  " + bold_prefix
        _font(r1, size=size, bold=True, color=color)
        r2 = p.add_run()
        r2.text = text
        _font(r2, size=size, bold=False, color=color)
    else:
        r = p.add_run()
        r.text = "\u2022  " + text
        _font(r, size=size, bold=False, color=color)


def _accent_line(slide, left, top, width=Inches(1.1), height=Pt(5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TERRACOTTA
    shape.line.fill.background()
    return shape


def _section_header(slide, label, title, top_start=None):
    y = top_start or M_TOP
    _add_textbox(slide, M_LEFT, y, CONTENT_W, Inches(0.4),
                 text=label, size=14, bold=True, color=TERRACOTTA)
    y += Inches(0.45)
    _add_textbox(slide, M_LEFT, y, CONTENT_W, Inches(0.6),
                 text=title, size=28, bold=True, color=NEAR_BLACK)
    y += Inches(0.65)
    _accent_line(slide, M_LEFT, y)
    y += Inches(0.35)
    return y


def _card_shape(slide, left, top, width, height, fill=WHITE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = RGBColor(232, 227, 214)
    card.line.width = Pt(1)
    return card


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_tag(slide, left, top, text, bg_color=TERRACOTTA, text_color=PARCHMENT):
    w = Inches(max(1.0, len(text) * 0.12 + 0.3))
    h = Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _font(r, size=10, bold=True, color=text_color)
    return w


def _icon_circle(slide, left, top, size, color, symbol):
    """Draw a colored circle with a Unicode symbol."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = symbol
    _font(r, size=int(size * 20), bold=True, color=PARCHMENT)
    return shape


# ---------------------------------------------------------------------------
# Slide 9A: Responsible AI Framework
# ---------------------------------------------------------------------------
def slide_ethics_framework(prs):
    """AI Ethics — Responsible AI Framework"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, PARCHMENT)
    y = _section_header(slide, "AI ETHICS", "Responsible AI by Design")

    # Intro text
    intro_box = _add_textbox(
        slide, M_LEFT, y, CONTENT_W, Inches(0.6),
        text="Cervantes embeds ethical safeguards at every layer — from content generation to student interaction to data handling.",
        size=15, color=BODY_TEXT
    )
    y += Inches(0.7)

    # 6 principle cards in 3x2 grid
    principles = [
        (MASTERY, "\u2693",  # anchor
         "Human Oversight",
         "Teachers review and approve all AI-generated content before students see it. "
         "AI proposes — educators decide."),
        (TERRACOTTA, "\u2696",  # scales
         "Fairness & Inclusion",
         "Socratic pushback evaluates reasoning quality, not language fluency or "
         "cultural background. No penalty for non-standard expression."),
        (WHEAT_GOLD, "\u26A0",  # warning
         "Transparency",
         "Students know they are interacting with AI. "
         "Teachers see exactly which AI models generated each scene and assessment."),
        (MASTERY, "\u2615",  # shield-like
         "Student Wellbeing",
         "Narrative format reduces exam anxiety. "
         "Socratic dialogue encourages growth, never punishes mistakes or uncertainty."),
        (CRITICAL, "\u2660",  # spade/lock
         "Data Minimisation",
         "Reasoning traces are treated as protected educational records. "
         "Only data necessary for assessment is collected and retained."),
        (WHEAT_GOLD, "\u2699",  # gear
         "Robustness & Safety",
         "Gemini outputs are validated against curriculum standards via CurricuLLM. "
         "Retry logic and fallbacks prevent hallucinated assessments."),
    ]

    card_w = Inches(3.7)
    card_h = Inches(2.0)
    gap_x = Inches(0.3)
    gap_y = Inches(0.2)

    for i, (icon_color, icon_sym, title, desc) in enumerate(principles):
        col = i % 3
        row = i // 3
        cx = M_LEFT + col * (card_w + gap_x)
        cy = y + row * (card_h + gap_y)

        card = _card_shape(slide, cx, cy, card_w, card_h)
        tf = card.text_frame
        tf.word_wrap = True

        # Icon + Title line
        p0 = tf.paragraphs[0]
        p0.space_after = Pt(6)
        r_icon = p0.add_run()
        r_icon.text = icon_sym + "  "
        _font(r_icon, size=16, bold=True, color=icon_color)
        r_title = p0.add_run()
        r_title.text = title
        _font(r_title, size=14, bold=True, color=NEAR_BLACK)

        _add_para(tf, desc, size=11, color=BODY_TEXT, space_after=Pt(2))


# ---------------------------------------------------------------------------
# Slide 9B: Data Privacy & Compliance
# ---------------------------------------------------------------------------
def slide_ethics_privacy(prs):
    """AI Ethics — Data Privacy & Student Protection"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, NEAR_BLACK)
    y = M_TOP

    # Header (light on dark)
    _add_textbox(slide, M_LEFT, y, CONTENT_W, Inches(0.4),
                 text="AI ETHICS", size=14, bold=True, color=WHEAT_GOLD)
    y += Inches(0.45)
    _add_textbox(slide, M_LEFT, y, CONTENT_W, Inches(0.6),
                 text="Data Privacy & Student Protection", size=28,
                 bold=True, color=PARCHMENT)
    y += Inches(0.65)
    _accent_line(slide, M_LEFT, y)
    y += Inches(0.4)

    # Left column — What We Protect
    left_card = _card_shape(slide, M_LEFT, y, Inches(5.5), Inches(4.0),
                            fill=DARK_BG)
    left_card.line.color.rgb = RGBColor(80, 75, 65)
    tf_left = left_card.text_frame
    tf_left.word_wrap = True

    p0 = tf_left.paragraphs[0]
    p0.space_after = Pt(10)
    r0 = p0.add_run()
    r0.text = "WHAT WE PROTECT"
    _font(r0, size=12, bold=True, color=TERRACOTTA)

    protections = [
        ("Reasoning Traces", " — classified as protected educational records, never shared with third parties"),
        ("Student Identity", " — randomised character assignments prevent cross-student data leakage"),
        ("Conversation Data", " — stored in Firestore with project-level access controls"),
        ("AI Model Inputs", " — student responses sent to Gemini API with zero-retention policy"),
        ("Assessment Outcomes", " — visible only to the assigned professor, not other students"),
    ]
    for bold_part, rest in protections:
        _add_bullet(tf_left, rest, size=12, color=PARCHMENT, bold_prefix=bold_part)

    # Right column — Compliance & Standards
    rx = Inches(7.0)
    right_card = _card_shape(slide, rx, y, Inches(5.2), Inches(4.0),
                             fill=DARK_BG)
    right_card.line.color.rgb = RGBColor(80, 75, 65)
    tf_right = right_card.text_frame
    tf_right.word_wrap = True

    p1 = tf_right.paragraphs[0]
    p1.space_after = Pt(10)
    r1 = p1.add_run()
    r1.text = "ALIGNMENT & COMMITMENTS"
    _font(r1, size=12, bold=True, color=WHEAT_GOLD)

    standards = [
        ("FERPA-Ready", " — architecture supports Family Educational Rights and Privacy Act compliance"),
        ("GDPR Principles", " — data minimisation, purpose limitation, and right to erasure supported"),
        ("Australian Privacy Act", " — handles student data in accordance with APP guidelines"),
        ("UNESCO AI Ethics", " — aligned with UNESCO Recommendation on Ethics of AI in Education"),
        ("Google AI Principles", " — Gemini usage follows Google's Responsible AI commitments"),
    ]
    for bold_part, rest in standards:
        _add_bullet(tf_right, rest, size=12, color=PARCHMENT, bold_prefix=bold_part)

    # Bottom banner
    banner_y = y + Inches(4.3)
    banner = _card_shape(slide, M_LEFT, banner_y, CONTENT_W, Inches(0.7),
                         fill=TERRACOTTA)
    banner.line.fill.background()
    tf_banner = banner.text_frame
    tf_banner.word_wrap = True
    try:
        tf_banner.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p_b = tf_banner.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    r_b = p_b.add_run()
    r_b.text = "AI generates and probes  \u2014  Educators control and evaluate  \u2014  Students own their data"
    _font(r_b, size=14, bold=True, color=PARCHMENT)


# ---------------------------------------------------------------------------
# Slide 9C: Bias Mitigation & Assessment Integrity
# ---------------------------------------------------------------------------
def slide_ethics_bias(prs):
    """AI Ethics — Bias Mitigation & Assessment Integrity"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, PARCHMENT)
    y = _section_header(slide, "AI ETHICS", "Bias Mitigation & Assessment Integrity")

    # Left — How We Mitigate Bias
    left_box = slide.shapes.add_textbox(M_LEFT, y, Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True

    p_h = tf.paragraphs[0]
    p_h.space_after = Pt(8)
    r_h = p_h.add_run()
    r_h.text = "HOW WE MITIGATE BIAS"
    _font(r_h, size=12, bold=True, color=TERRACOTTA)

    mitigations = [
        ("Curriculum Grounding",
         " — CurricuLLM anchors all content to official standards, preventing AI hallucination or off-topic assessment"),
        ("Character Randomisation",
         " — students receive different character names, personalities, and sprite variants to prevent stereotyping"),
        ("Reasoning-First Evaluation",
         " — signal extraction judges logical reasoning quality, not grammar, vocabulary, or writing style"),
        ("Socratic Neutrality",
         " — pushback dialogue challenges all students equally, probing misconceptions without assuming deficit"),
        ("Multi-Path Assessment",
         " — branching dialogue allows multiple valid reasoning approaches, respecting diverse thinking styles"),
        ("Teacher Override",
         " — professors can edit any AI-generated scene, question, or misconception target before publishing"),
    ]
    for bold_part, rest in mitigations:
        _add_bullet(tf, rest, size=12, color=BODY_TEXT, bold_prefix=bold_part)

    # Right — Trust Architecture diagram
    rx = Inches(7.0)
    card = _card_shape(slide, rx, y, Inches(5.2), Inches(4.5))
    tf2 = card.text_frame
    tf2.word_wrap = True
    try:
        tf2.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass

    p_t = tf2.paragraphs[0]
    p_t.space_after = Pt(12)
    r_t = p_t.add_run()
    r_t.text = "TRUST ARCHITECTURE"
    _font(r_t, size=12, bold=True, color=WHEAT_GOLD)

    trust_layers = [
        ("\u2460  Curriculum Layer",
         "CurricuLLM validates learning objectives against national standards"),
        ("\u2461  Generation Layer",
         "Gemini produces content within curriculum-defined boundaries"),
        ("\u2462  Review Layer",
         "Professor reviews, edits, and approves before publication"),
        ("\u2463  Interaction Layer",
         "Socratic dialogue probes reasoning without leading or biasing"),
        ("\u2464  Evaluation Layer",
         "Signal extraction uses rubric-aligned criteria, not AI opinion"),
        ("\u2465  Audit Layer",
         "Full reasoning traces stored for transparency and appeals"),
    ]
    for title, desc in trust_layers:
        _add_para(tf2, title, size=13, bold=True, color=NEAR_BLACK,
                  space_before=Pt(6), space_after=Pt(2))
        _add_para(tf2, desc, size=11, color=BODY_TEXT, space_after=Pt(4))


# ---------------------------------------------------------------------------
# Main — insert ethics slides into the Google PPTX
# ---------------------------------------------------------------------------
def main():
    prs = Presentation(str(INPUT_FILE))

    # Add the three ethics slides (appended at end)
    slide_ethics_framework(prs)
    slide_ethics_privacy(prs)
    slide_ethics_bias(prs)

    # Reorder: move new slides (last 3) to after slide 8 (index 7)
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldIdLst = prs.element.find('.//p:sldIdLst', nsmap)

    if sldIdLst is not None:
        sldIds = list(sldIdLst)
        total = len(sldIds)

        # Original: [0..9] = slides 1-10, then [10,11,12] = 3 new ethics slides
        # Target: [0..7, 10, 11, 12, 8, 9] = slides 1-8, ethics x3, then 9-10
        if total >= 13:
            new_order = list(range(8)) + [10, 11, 12] + [8, 9]
            reordered = [sldIds[i] for i in new_order]

            for child in list(sldIdLst):
                sldIdLst.remove(child)
            for elem in reordered:
                sldIdLst.append(elem)

    prs.save(str(OUTPUT_FILE))
    print(f"Saved: {OUTPUT_FILE}")
    print(f"Total slides: {total}")

    # Verify slide order
    prs2 = Presentation(str(OUTPUT_FILE))
    for i, slide in enumerate(prs2.slides, 1):
        first_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        first_text = t
                        break
            if first_text:
                break
        print(f"  Slide {i:2d}: {first_text[:60]}")


if __name__ == "__main__":
    main()
