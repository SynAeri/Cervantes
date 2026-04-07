#!/usr/bin/env python3
"""
generate_pptx.py — Generates cervantes-deck.pptx from the La Mancha Parchment design system.

Content extracted from index.html (Cambridge EdTech x Hack 2026 slide deck).
Requires: python-pptx >= 1.0.2
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Design tokens — La Mancha Parchment
# ---------------------------------------------------------------------------
TERRACOTTA = RGBColor(200, 90, 50)
WHEAT_GOLD = RGBColor(196, 162, 69)
PARCHMENT = RGBColor(244, 241, 234)
NEAR_BLACK = RGBColor(30, 28, 24)
WHITE = RGBColor(255, 255, 255)
BODY_TEXT = RGBColor(58, 55, 48)       # #3a3730
MUTED_TEXT = RGBColor(107, 101, 90)    # #6b655a
MASTERY = RGBColor(59, 130, 126)       # #3B827E
MISCONCEPTION = RGBColor(212, 163, 71) # #D4A347
CRITICAL = RGBColor(158, 59, 59)       # #9E3B3B
SERVICES_GREEN = RGBColor(107, 143, 107)  # #6b8f6b
DARK_BG = RGBColor(42, 39, 32)        # #2a2720

FONT_NAME = "Plus Jakarta Sans"
FALLBACK_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Margins
M_LEFT = Inches(0.8)
M_TOP = Inches(0.6)
CONTENT_W = Inches(11.7)

OUTPUT = Path(__file__).parent / "cervantes-deck.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _font(run, size=18, bold=False, color=NEAR_BLACK, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _add_textbox(slide, left, top, width, height, text="", size=18,
                 bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, font_name=FONT_NAME):
    """Add a simple textbox and return the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _font(run, size=size, bold=bold, color=color, name=font_name)
    return txBox


def _add_para(text_frame, text, size=18, bold=False, color=NEAR_BLACK,
              space_before=Pt(0), space_after=Pt(4), align=PP_ALIGN.LEFT,
              font_name=FONT_NAME):
    """Append a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    _font(run, size=size, bold=bold, color=color, name=font_name)
    return p


def _add_bullet(text_frame, text, size=18, color=BODY_TEXT, indent=0,
                bold_prefix="", font_name=FONT_NAME):
    """Add a bullet point (with optional bold prefix)."""
    p = text_frame.add_paragraph()
    p.space_before = Pt(2)
    p.space_after = Pt(4)
    p.level = indent
    # Bullet character
    if bold_prefix:
        r1 = p.add_run()
        r1.text = "\u2022  " + bold_prefix
        _font(r1, size=size, bold=True, color=color, name=font_name)
        r2 = p.add_run()
        r2.text = text
        _font(r2, size=size, bold=False, color=color, name=font_name)
    else:
        r = p.add_run()
        r.text = "\u2022  " + text
        _font(r, size=size, bold=False, color=color, name=font_name)


def _accent_line(slide, left, top, width=Inches(1.1), height=Pt(5)):
    """Draw the terracotta-to-gold accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TERRACOTTA
    shape.line.fill.background()
    return shape


def _section_header(slide, label, title, top_start=None):
    """Add the standard section-label + h2 + accent-line block. Returns y after."""
    y = top_start or M_TOP
    _add_textbox(slide, M_LEFT, y, CONTENT_W, Inches(0.4),
                 text=label, size=14, bold=True, color=TERRACOTTA)
    y += Inches(0.45)
    _add_textbox(slide, M_LEFT, y, CONTENT_W, Inches(0.6),
                 text=title, size=28, bold=True, color=NEAR_BLACK)
    y += Inches(0.65)
    _accent_line(slide, M_LEFT, y)
    y += Inches(0.35)
    return y


def _video_placeholder(slide, left, top, width, height, label_text):
    """Draw a dark rectangle with centered label — video placeholder."""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK_BG
    rect.line.color.rgb = WHEAT_GOLD
    rect.line.width = Pt(1.5)
    rect.line.dash_style = 4  # dash

    tf = rect.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    # Play icon symbol
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = "\u25B6"
    _font(r0, size=36, bold=True, color=TERRACOTTA)

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = label_text
    _font(r1, size=12, bold=True, color=WHEAT_GOLD)
    return rect


def _card_shape(slide, left, top, width, height):
    """White rounded-rectangle card."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(232, 227, 214)  # parchment-dark
    card.line.width = Pt(1)
    return card


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_tag(slide, left, top, text, bg_color=TERRACOTTA, text_color=PARCHMENT):
    """Small tag pill."""
    w = Inches(max(1.0, len(text) * 0.12 + 0.3))
    h = Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _font(r, size=10, bold=True, color=text_color)
    return w


def _numbered_step(text_frame, num, title, desc, size=14):
    """Add a numbered step row to a text frame."""
    p = text_frame.add_paragraph()
    p.space_before = Pt(6)
    p.space_after = Pt(2)
    r_num = p.add_run()
    r_num.text = f"{num}. "
    _font(r_num, size=size, bold=True, color=TERRACOTTA)
    r_title = p.add_run()
    r_title.text = title
    _font(r_title, size=size, bold=True, color=NEAR_BLACK)

    p2 = text_frame.add_paragraph()
    p2.space_before = Pt(0)
    p2.space_after = Pt(6)
    r_desc = p2.add_run()
    r_desc.text = "     " + desc
    _font(r_desc, size=12, bold=False, color=BODY_TEXT)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def slide_01_title(prs):
    """Title slide: CERVANTES"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, PARCHMENT)

    # Logo mark
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(6.0), Inches(1.0), Inches(1.1), Inches(1.1))
    logo.fill.solid()
    logo.fill.fore_color.rgb = TERRACOTTA
    logo.line.fill.background()
    tf = logo.text_frame
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "C"
    _font(r, size=40, bold=True, color=PARCHMENT)

    # Title
    _add_textbox(slide, Inches(1.5), Inches(2.4), Inches(10.3), Inches(1.2),
                 text="CERVANTES", size=54, bold=True, color=NEAR_BLACK,
                 align=PP_ALIGN.CENTER)

    # Accent line centered
    _accent_line(slide, Inches(5.8), Inches(3.65), width=Inches(1.7))

    # Subtitle
    _add_textbox(slide, Inches(2.0), Inches(3.95), Inches(9.3), Inches(0.6),
                 text="Narrative-Driven Assessment for Higher Education",
                 size=20, bold=False, color=TERRACOTTA, align=PP_ALIGN.CENTER)

    # Event label
    _add_textbox(slide, Inches(2.0), Inches(4.7), Inches(9.3), Inches(0.4),
                 text="Cambridge EdTech x Hack 2026  \u2014  EduX Challenge 1",
                 size=14, bold=False, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

    # Team tags
    x = Inches(5.0)
    y = Inches(5.5)
    for name in ["Jordan", "Liu", "Adrian"]:
        w = _add_tag(slide, x, y, name)
        x += w + Inches(0.15)


def slide_02_problem(prs):
    """The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PARCHMENT)
    y = _section_header(slide, "THE PROBLEM", "Traditional Assessments Are Broken")

    # Left column — bullet list
    left_box = slide.shapes.add_textbox(M_LEFT, y, Inches(5.5), Inches(3.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    _add_bullet(tf, " measure recall under pressure, not genuine understanding",
                size=15, bold_prefix="High-stakes exams")
    _add_bullet(tf, " reduce rich learning to multiple choice and short answers",
                size=15, bold_prefix="Rigid formats")
    _add_bullet(tf, " \u2014 educators see scores, never thought processes",
                size=15, bold_prefix="No reasoning visibility")
    _add_bullet(tf, " distorts performance and undermines learning outcomes",
                size=15, bold_prefix="Student anxiety")

    # Right column — "The Gap" card
    rx = Inches(7.0)
    card = _card_shape(slide, rx, y, Inches(5.2), Inches(3.2))
    tf2 = card.text_frame
    tf2.word_wrap = True
    p0 = tf2.paragraphs[0]
    p0.space_after = Pt(8)
    r0 = p0.add_run()
    r0.text = "THE GAP"
    _font(r0, size=12, bold=True, color=TERRACOTTA)

    _add_para(tf2, "Universities need assessments that reveal how students think \u2014 "
              "not just what they remember.", size=15, color=BODY_TEXT,
              space_after=Pt(10))
    _add_para(tf2, "Current tools offer no middle ground between hand-graded "
              "essays and automated quizzes.", size=15, color=BODY_TEXT)


def slide_03_solution(prs):
    """Our Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PARCHMENT)
    y = _section_header(slide, "OUR SOLUTION", "Assessment Through Interactive Narrative")

    # Left column — description + steps
    left_box = slide.shapes.add_textbox(M_LEFT, y, Inches(5.5), Inches(4.2))
    tf = left_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.space_after = Pt(12)
    r0 = p0.add_run()
    r0.text = ("Cervantes transforms assessment into Visual Novel experiences \u2014 "
               "curriculum-aligned, AI-generated scenarios where students demonstrate "
               "understanding through dialogue and decision-making.")
    _font(r0, size=15, color=BODY_TEXT)

    _numbered_step(tf, 1, "Narrative Immersion",
                   "Students engage with characters and scenarios, not question sheets")
    _numbered_step(tf, 2, "Socratic Questioning",
                   "AI-driven dialogue probes understanding at natural checkpoints")
    _numbered_step(tf, 3, "Reasoning Traces",
                   "Every decision path is captured, giving educators deep insight")

    # Right column — card with summary
    rx = Inches(7.0)
    card = _card_shape(slide, rx, y, Inches(5.2), Inches(2.8))
    tf2 = card.text_frame
    tf2.word_wrap = True
    try:
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)
    r = p.add_run()
    r.text = "Textbook \u2192 Interactive Story \u2192 Rich Analytics"
    _font(r, size=16, bold=True, color=NEAR_BLACK)

    _add_para(tf2, "Textbook content becomes interactive story, "
              "producing rich analytics", size=14, color=BODY_TEXT,
              align=PP_ALIGN.CENTER, space_after=Pt(12))

    # Tags
    tag_y = y + Inches(3.1)
    tx = rx + Inches(0.5)
    for tag_text in ["Curriculum-Aligned", "AI-Generated", "Real-Time"]:
        w = _add_tag(slide, tx, tag_y, tag_text, bg_color=WHEAT_GOLD,
                     text_color=NEAR_BLACK)
        tx += w + Inches(0.15)


def slide_04_teacher_flow(prs):
    """Teacher Flow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PARCHMENT)
    y = _section_header(slide, "TEACHER FLOW", "From Curriculum to Classroom in Minutes")

    # Left column — 4 steps
    left_box = slide.shapes.add_textbox(M_LEFT, y, Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True

    _numbered_step(tf, 1, "Upload Materials",
                   "Subject, module, concept targets, rubric focus (PDF / DOCX / TXT)")
    _numbered_step(tf, 2, "CurricuLLM Alignment",
                   "Auto-maps content to Australian curriculum standards and rubric structure")
    _numbered_step(tf, 3, "Narrative Generation",
                   "Gemini builds scenes, characters, Socratic checkpoints, and branching dialogue")
    _numbered_step(tf, 4, "Review & Publish",
                   "Approve the arc, assign to classes, and monitor live progress")

    # Right column — video placeholder
    _video_placeholder(slide, Inches(7.0), y, Inches(5.2), Inches(3.5),
                       "VIDEO OF TEACHER DASHBOARD DEMO GOES HERE")


def slide_05_student(prs):
    """Student Experience — dark background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NEAR_BLACK)
    y = M_TOP

    # Section header (light text on dark)
    _add_textbox(slide, M_LEFT, y, CONTENT_W, Inches(0.4),
                 text="STUDENT EXPERIENCE", size=14, bold=True, color=WHEAT_GOLD)
    y += Inches(0.45)
    _add_textbox(slide, M_LEFT, y, CONTENT_W, Inches(0.6),
                 text="Learn by Living the Story", size=28, bold=True, color=PARCHMENT)
    y += Inches(0.65)
    _accent_line(slide, M_LEFT, y)
    y += Inches(0.35)

    # Left column — video placeholder
    _video_placeholder(slide, M_LEFT, y, Inches(5.5), Inches(3.5),
                       "VIDEO OF STUDENT VN EXPERIENCE GOES HERE")

    # Right column — 3 cards
    rx = Inches(7.0)
    card_h = Inches(1.05)
    card_gap = Inches(0.15)
    cards_data = [
        ("Narrative Immersion",
         "Characters, scenes, and branching dialogue make assessment feel "
         "like gameplay \u2014 accessible via PWA on any device"),
        ("Socratic Checkpoints",
         "AI-driven questions appear at natural story beats, probing "
         "understanding without breaking immersion"),
        ("Reasoning Journaling",
         "Every choice, response, and dialogue path is captured in "
         "arc_journals and reasoning_traces"),
    ]
    cy = y
    for title, desc in cards_data:
        card = _card_shape(slide, rx, cy, Inches(5.2), card_h)
        tf = card.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.space_after = Pt(4)
        r0 = p0.add_run()
        r0.text = title
        _font(r0, size=11, bold=True, color=TERRACOTTA)
        _add_para(tf, desc, size=11, color=BODY_TEXT, space_after=Pt(2))
        cy += card_h + card_gap


def slide_06_intelligence(prs):
    """Intelligence Layer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PARCHMENT)
    y = _section_header(slide, "INTELLIGENCE LAYER",
                        "AI That Understands Curriculum and Cognition")

    # 3 cards in a row
    card_w = Inches(3.6)
    card_h = Inches(2.8)
    gap = Inches(0.45)
    cards = [
        ("CurricuLLM",
         "Grounds all assessment content in Australian curriculum standards",
         ["Concept target extraction", "Rubric structure alignment",
          "Misconception identification"]),
        ("Gemini 2.5 Flash",
         "Powers narrative generation and real-time dialogue",
         ["Scene & character generation", "Socratic questioning engine",
          "Adaptive branching logic"]),
        ("Reasoning Graphs",
         "Captures how students think, not just what they answer",
         ["Decision path tracking", "Misconception flagging",
          "Per-student trace views"]),
    ]
    cx = M_LEFT
    for title, desc, bullets in cards:
        card = _card_shape(slide, cx, y, card_w, card_h)
        tf = card.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.space_after = Pt(6)
        r0 = p0.add_run()
        r0.text = title
        _font(r0, size=12, bold=True, color=TERRACOTTA)
        _add_para(tf, desc, size=12, color=BODY_TEXT, space_after=Pt(8))
        for b in bullets:
            _add_bullet(tf, b, size=11, color=BODY_TEXT)
        cx += card_w + gap

    # Video placeholder below cards
    vp_y = y + card_h + Inches(0.3)
    _video_placeholder(slide, Inches(3.5), vp_y, Inches(6.3), Inches(1.8),
                       "VIDEO OF REASONING GRAPH DEMO GOES HERE")


def slide_07_architecture(prs):
    """Architecture — 3-column diagram with connectors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PARCHMENT)
    y = _section_header(slide, "ARCHITECTURE", "Full-Stack on Google Cloud")

    # Three architecture layers
    col_w = Inches(3.3)
    col_h = Inches(3.6)
    arrow_w = Inches(0.9)

    layers = [
        ("FRONTEND", WHEAT_GOLD, NEAR_BLACK,
         ["Next.js 16", "Teacher Dashboard", "Student PWA", "Vercel Hosting"]),
        ("BACKEND", TERRACOTTA, PARCHMENT,
         ["FastAPI", "Auth + User Mgmt", "Arc Generation",
          "Dialogue Engine", "Google Cloud Run"]),
        ("SERVICES", SERVICES_GREEN, WHITE,
         ["Firebase Auth + Firestore", "Gemini 2.5 Flash",
          "CurricuLLM Standards", "Google Cloud Platform"]),
    ]

    x = M_LEFT
    layer_shapes = []
    for i, (label, badge_bg, badge_fg, items) in enumerate(layers):
        # Column box
        col = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, y, col_w, col_h)
        col.fill.solid()
        col.fill.fore_color.rgb = WHITE
        col.line.color.rgb = badge_bg
        col.line.width = Pt(2)

        tf = col.text_frame
        tf.word_wrap = True

        # Badge label
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(10)
        r0 = p0.add_run()
        r0.text = label
        _font(r0, size=10, bold=True, color=badge_fg)

        # Add a small colored badge behind the title
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x + Inches(0.8), y + Inches(0.15),
                                       Inches(1.7), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_bg
        badge.line.fill.background()
        btf = badge.text_frame
        try:
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = label
        _font(br, size=10, bold=True, color=badge_fg)

        # Items
        for idx, item in enumerate(items):
            is_last = (idx == len(items) - 1)
            clr = MUTED_TEXT if is_last else BODY_TEXT
            bld = (idx == 0)
            _add_para(tf, item, size=13, bold=bld, color=clr,
                      align=PP_ALIGN.CENTER, space_before=Pt(6),
                      space_after=Pt(6))

        layer_shapes.append((x, y))
        x += col_w

        # Arrow between columns
        if i < 2:
            arrow_x = x + Inches(0.05)
            # Arrow label box
            lbl_top = ["REST / API", "SDK / APIs"]
            _add_textbox(slide, arrow_x, y + Inches(1.4), arrow_w, Inches(0.4),
                         text=lbl_top[i], size=10, bold=True, color=TERRACOTTA,
                         align=PP_ALIGN.CENTER)
            # Arrow line shape
            arrow_shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x + Inches(0.1),
                y + Inches(1.8), Inches(0.7), Inches(0.25))
            arrow_shape.fill.solid()
            arrow_shape.fill.fore_color.rgb = TERRACOTTA
            arrow_shape.line.fill.background()

            x += arrow_w


def slide_08_tech_highlights(prs):
    """Technical Highlights — 2x3 grid (actually 5 cards in the HTML)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PARCHMENT)
    y = _section_header(slide, "TECH HIGHLIGHTS", "What Makes It Work")

    # Row 1: 3 cards
    card_w = Inches(3.6)
    card_h = Inches(1.8)
    gap = Inches(0.45)

    row1 = [
        ("Curriculum Grounding",
         "CurricuLLM ensures every generated scenario maps to real curriculum "
         "standards \u2014 no hallucinated assessments"),
        ("Real-Time Traces",
         "Firestore listeners push live reasoning data to teacher dashboards "
         "as students progress through scenarios"),
        ("Adaptive Dialogue",
         "Gemini-powered Socratic engine adjusts questioning depth based on "
         "student responses in real-time"),
    ]
    cx = M_LEFT
    for title, desc in row1:
        card = _card_shape(slide, cx, y, card_w, card_h)
        tf = card.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.space_after = Pt(6)
        r0 = p0.add_run()
        r0.text = title
        _font(r0, size=12, bold=True, color=TERRACOTTA)
        _add_para(tf, desc, size=12, color=BODY_TEXT)
        cx += card_w + gap

    # Row 2: 2 wider cards (+ 1 filler to make a 2x3 look)
    y2 = y + card_h + Inches(0.3)
    wide_w = Inches(5.5)
    row2 = [
        ("PWA-First Mobile",
         "Students access VN assessments on any device via Progressive Web "
         "App \u2014 no app store, no installation"),
        ("3 Core Endpoints",
         "Arc Setup (teacher upload + CurricuLLM), Narrative Generation "
         "(Gemini scenes), Dialogue (real-time student interaction)"),
    ]
    cx = M_LEFT
    for title, desc in row2:
        card = _card_shape(slide, cx, y2, wide_w, card_h)
        tf = card.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.space_after = Pt(6)
        r0 = p0.add_run()
        r0.text = title
        _font(r0, size=12, bold=True, color=TERRACOTTA)
        _add_para(tf, desc, size=12, color=BODY_TEXT)
        cx += wide_w + gap

    # Add a 6th card: Docker (from task spec)
    card6 = _card_shape(slide, cx - gap - Inches(0.1), y2, Inches(1.4), card_h)
    tf6 = card6.text_frame
    tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    r6 = p6.add_run()
    r6.text = "Docker"
    _font(r6, size=12, bold=True, color=TERRACOTTA)
    _add_para(tf6, "Containerised deployment", size=10, color=BODY_TEXT)


def slide_09_impact(prs):
    """Impact & Future"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PARCHMENT)
    y = _section_header(slide, "IMPACT & FUTURE", "Redefining Assessment at Scale")

    # Left column — stats + what changes
    # Stat boxes
    stat_w = Inches(2.4)
    stat_h = Inches(1.3)
    for i, (val, label) in enumerate([("0%", "Exam Anxiety"),
                                       ("100%", "Reasoning Visibility")]):
        sx = M_LEFT + i * (stat_w + Inches(0.3))
        card = _card_shape(slide, sx, y, stat_w, stat_h)
        tf = card.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(4)
        r0 = p0.add_run()
        r0.text = val
        _font(r0, size=36, bold=True, color=TERRACOTTA)
        _add_para(tf, label, size=13, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

    # "What Changes" bullets
    wc_y = y + stat_h + Inches(0.3)
    wc_box = slide.shapes.add_textbox(M_LEFT, wc_y, Inches(5.2), Inches(2.0))
    tf = wc_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "WHAT CHANGES"
    _font(r0, size=12, bold=True, color=TERRACOTTA)
    _add_bullet(tf, " not pressure",
                size=14, bold_prefix="Students assessed through engagement,")
    _add_bullet(tf, " not just grades",
                size=14, bold_prefix="Educators see thought processes,")
    _add_bullet(tf, " not manual",
                size=14, bold_prefix="Curriculum alignment is automated,")

    # Right column — Roadmap
    rx = Inches(7.0)
    _add_textbox(slide, rx, y, Inches(5.2), Inches(0.4),
                 text="ROADMAP", size=12, bold=True, color=TERRACOTTA)

    road_card = _card_shape(slide, rx, y + Inches(0.45), Inches(5.2), Inches(3.2))
    tf2 = road_card.text_frame
    tf2.word_wrap = True
    _add_bullet(tf2, " narrative support", size=14, bold_prefix="Multi-language")
    _add_bullet(tf2, " (Moodle, Canvas, Blackboard)", size=14,
                bold_prefix="LMS integration")
    _add_bullet(tf2, " multiplayer VN scenarios", size=14,
                bold_prefix="Collaborative")
    _add_bullet(tf2, " beyond Australian standards", size=14,
                bold_prefix="Global curricula")
    _add_bullet(tf2, " via Gemini TTS/STT", size=14,
                bold_prefix="Voice interaction")


def slide_10_thankyou(prs):
    """Thank You — dark slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NEAR_BLACK)

    # Title
    _add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.0),
                 text="Thank You", size=48, bold=True, color=PARCHMENT,
                 align=PP_ALIGN.CENTER)

    # Accent line
    _accent_line(slide, Inches(5.8), Inches(2.35), width=Inches(1.7))

    # Tagline
    _add_textbox(slide, Inches(2.0), Inches(2.7), Inches(9.3), Inches(0.5),
                 text="Redefining assessment, one story at a time.",
                 size=18, bold=False, color=RGBColor(200, 196, 186),
                 align=PP_ALIGN.CENTER)

    # Team grid
    team = [("Jordan", "Frontend + Narrative Arc"),
            ("Liu", "Backend + Cloud"),
            ("Adrian", "Assets + Infrastructure")]
    tx = Inches(2.0)
    for name, role in team:
        _add_textbox(slide, tx, Inches(3.6), Inches(3.0), Inches(0.4),
                     text=name, size=18, bold=True, color=PARCHMENT,
                     align=PP_ALIGN.CENTER)
        _add_textbox(slide, tx, Inches(4.0), Inches(3.0), Inches(0.35),
                     text=role, size=13, bold=False, color=WHEAT_GOLD,
                     align=PP_ALIGN.CENTER)
        tx += Inches(3.3)

    # GitHub
    _add_textbox(slide, Inches(2.0), Inches(5.0), Inches(9.3), Inches(0.4),
                 text="GitHub   github.com/SynAeri/IncubedHackathon",
                 size=14, bold=False, color=RGBColor(200, 196, 186),
                 align=PP_ALIGN.CENTER)

    # Tech tags
    tags = [("Next.js 16", TERRACOTTA, PARCHMENT),
            ("FastAPI", TERRACOTTA, PARCHMENT),
            ("Firebase", TERRACOTTA, PARCHMENT),
            ("Gemini 2.5 Flash", TERRACOTTA, PARCHMENT),
            ("CurricuLLM", WHEAT_GOLD, NEAR_BLACK)]
    tx = Inches(3.2)
    for text, bg, fg in tags:
        w = _add_tag(slide, tx, Inches(5.7), text, bg_color=bg, text_color=fg)
        tx += w + Inches(0.15)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_teacher_flow(prs)
    slide_05_student(prs)
    slide_06_intelligence(prs)
    slide_07_architecture(prs)
    slide_08_tech_highlights(prs)
    slide_09_impact(prs)
    slide_10_thankyou(prs)

    prs.save(str(OUTPUT))
    print(f"Saved {OUTPUT}  ({OUTPUT.stat().st_size / 1024:.1f} KB, "
          f"{len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
